from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 创建演示文稿对象
prs = Presentation()

# 定义颜色
BRONZE = RGBColor(140, 120, 83)  # 古铜色
DARK_GREY = RGBColor(50, 50, 50)
BG_COLOR = RGBColor(250, 250, 245) # 仿纸张米色

def set_slide_background(slide):
    """设置简单的背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # 标题幻灯片
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = BRONZE
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.name = "微软雅黑"
    
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.color.rgb = DARK_GREY
    subtitle.text_frame.paragraphs[0].font.name = "微软雅黑"

def add_content_slide(prs, title_text, content_text_list):
    slide_layout = prs.slide_layouts[1] # 标题+内容
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    # 设置标题
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = BRONZE
    title.text_frame.paragraphs[0].font.name = "微软雅黑"
    title.text_frame.paragraphs[0].font.bold = True
    
    # 设置正文
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear() # 清除默认文本
    
    for item in content_text_list:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GREY
        p.font.name = "微软雅黑"
        p.space_after = Pt(14)

def add_table_slide(prs, title_text, data):
    slide_layout = prs.slide_layouts[5] # 仅标题
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = BRONZE
    title.text_frame.paragraphs[0].font.name = "微软雅黑"
    
    # 添加表格
    rows = len(data)
    cols = len(data[0])
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(0.8 * rows)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.text_frame.paragraphs[0].font.name = "微软雅黑"
            
# --- 生成内容 ---

# 1. 封面
add_title_slide(prs, "青铜器上的东方纹样", "小学美术“综合·探索”单元教学设计汇报\n\n深圳市景贝小学  李婷婷")

# 2. 单元构成
data_unit = [
    ["单元主题", "青铜器上的东方纹样", "单元课时", "4课时"],
    ["任务类型", "综合·探索", "授课年级", "三年级"],
    ["教材版本", "岭南新版《艺术》", "设计者", "李婷婷"]
]
add_table_slide(prs, "一、单元构成与背景", data_unit)

# 3. 单元架构
content_structure = [
    "第一课：感知 —— 《青铜密码》",
    "   - 关键词：盲盒、分类、简笔画",
    "   - 目标：认识礼器兵器，归纳纹样规律",
    "",
    "第二课：探究 —— 《纹样探秘》",
    "   - 关键词：拼图、规律、设计",
    "   - 目标：掌握对称/重复/渐变，器物适配",
    "",
    "第三课：实践 —— 《立体创作》",
    "   - 关键词：锡箔纸、技法、融合",
    "   - 目标：掌握压印/剪贴/揉搓，表现立体层次",
    "",
    "第四课：展示 —— 《灵思绘韵》",
    "   - 关键词：策展、解说、互评",
    "   - 目标：文化解码，综合评价"
]
add_content_slide(prs, "二、单元整体架构", content_structure)

# 4. 目标一致性
content_goals = [
    "【知识与技能】",
    "了解青铜历史背景，掌握对称、重复、渐变规律。创新纹样设计。",
    "",
    "【过程与方法】",
    "观察分析 -> 锡箔纸技法实践 -> 仿青铜制作 -> 多元评价。",
    "",
    "【情感态度价值观】",
    "增强文化认同感，激发传统元素融入现代生活的创新意识。"
]
add_content_slide(prs, "三、单元任务与目标一致性", content_goals)

# 5. 第一课设计
content_l1 = [
    "课题：《青铜密码——纹样里的东方故事》",
    "",
    "1. 教学前：前置知识问卷，兴趣调查。",
    "2. 教学中（核心环节）：",
    "   - 考古队长引入，开“纹样盲盒”。",
    "   - 磁贴板分类游戏（礼器 vs 兵器）。",
    "   - 传统纹样 + 现代物品 = 创意简笔画。",
    "3. 测评工具：",
    "   - 《观察记录表》、创意画评分表。"
]
add_content_slide(prs, "四、分课设计：第1课", content_l1)

# 6. 第二课设计
content_l2 = [
    "课题：《纹样探秘——设计里的规律之美》",
    "",
    "1. 核心任务：",
    "   - 规律解码：拼图还原纹样（对称/重复/渐变）。",
    "   - 器物适配：圆形鼎、方形尊的纹样选择。",
    "2. 深度探究：",
    "   - 使用“设计检查清单”进行自我监控。",
    "   - 用“因为...所以...”句式解释设计理由。",
    "3. 测评工具：",
    "   - 《拼图实验评价表》、《合作设计评价表》。"
]
add_content_slide(prs, "四、分课设计：第2课", content_l2)

# 7. 第三课设计 (重点)
content_l3 = [
    "课题：《立体创作——锡箔纸上的青铜新语》",
    "",
    "1. 材料探秘：",
    "   - 视频学习：揉皱、雕刻、多层粘贴技法。",
    "2. 创意实践：",
    "   - 盲盒挑战：传统纹样 + 现代科技感。",
    "   - 锡箔纸仿青铜工艺制作。",
    "3. 展示评价：",
    "   - 光影剧场：投影展示，故事接龙。",
    "4. 测评工具：",
    "   - 作品评价量规、创意星投票。"
]
add_content_slide(prs, "四、分课设计：第3课（重点）", content_l3)

# 8. 第四课设计
content_l4 = [
    "课题：《灵思绘韵——青铜纹样策展幻旅》",
    "",
    "1. 策展训练：",
    "   - 展区布置：文化解码、设计亮点讲解。",
    "2. 观众互动：",
    "   - 改良盲盒、便签墙反馈、问答闯关。",
    "3. 动态改良：",
    "   - 根据观众建议绘制改良草图。",
    "4. 测评工具：",
    "   - “三明治评价法”、观众投票、策展人大赛。"
]
add_content_slide(prs, "四、分课设计：第4课", content_l4)

# 9. 实施亮点 - 导入
content_highlight_1 = [
    "聚焦第3课：《立体创作》实施细节",
    "",
    "【教学前评价渗透】",
    "问题情境：如果青铜器上的纹样会说话，它想传递什么？",
    "",
    "【环节一：纹样挑战】",
    "- 观察：饕餮纹鼎、宴乐纹壶细节。",
    "- 动作：分类（几何/动物/人物），提取关键词。",
    "- 猜测：纹样背后的文化寓意。"
]
add_content_slide(prs, "五、课程实施：情境与导入", content_highlight_1)

# 10. 实施亮点 - 过程
content_highlight_2 = [
    "【环节二：创意实践——锡箔纸上的“青铜重生”】",
    "",
    "技法四步走：",
    "1. 固定：美纹纸固定画稿与锡箔纸。",
    "2. 压痕：沿线条轻压，形成浮雕效果。",
    "3. 肌理：中心点打圈，做螺旋纹理增强层次。",
    "4. 做旧：涂仿青铜色颜料，模拟氧化效果。",
    "",
    "设计要求：构图饱满，线条流畅，结合现代元素（如手机壳）。"
]
add_content_slide(prs, "五、课程实施：技法与实践", content_highlight_2)

# 11. 实施亮点 - 评价
content_highlight_3 = [
    "【环节三：多元评价——让作品“会说话”】",
    "- 自评：《我的青铜纹样设计单》（谈理由、难点）。",
    "- 互评：使用话术“我发现……我建议……”。",
    "- 师评：颁发“青铜小匠人”勋章，点评线条动态。",
    "",
    "【环节四：展示升华】",
    "- 布置“微型青铜博物馆”。",
    "- 邀请校长/家长参观，颁发证书，增强仪式感。"
]
add_content_slide(prs, "五、课程实施：评价与展示", content_highlight_3)

# 12. 育人价值
content_value = [
    "1. 文化传承与创新",
    "   - 从“看历史”到“做历史”，在指尖感受文明的重量。",
    "   - 打破时空界限，将传统纹样应用于现代生活。",
    "",
    "2. 核心素养落地",
    "   - 审美感知：感受青铜纹样的秩序美与形式美。",
    "   - 艺术表现：掌握锡箔纸浮雕工艺，大胆表达。",
    "   - 文化理解：理解青铜礼乐文化，树立民族自信。"
]
add_content_slide(prs, "六、育人价值与课程反思", content_value)

# 13. 结束页
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
title = slide.shapes.title
title.text = "感谢聆听 敬请指正"
title.text_frame.paragraphs[0].font.color.rgb = BRONZE

# 保存文件
file_name = "青铜器上的东方纹样_汇报.pptx"
prs.save(file_name)
print(f"PPT文件已生成：{file_name}")